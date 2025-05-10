from fastapi import APIRouter, UploadFile, File, HTTPException, Form
from fastapi.responses import JSONResponse
from services.deepseek_api import generar_examen_deepseek, generar_respuestas_deepseek
import asyncio
import aiofiles
import os

from docx import Document
import fitz  # PyMuPDF
from pptx import Presentation

router = APIRouter()

def extraer_texto_docx(ruta):
    doc = Document(ruta)
    texto = "\n".join([p.text for p in doc.paragraphs])
    return texto

def extraer_texto_pdf(ruta):
    doc = fitz.open(ruta)
    texto = ""
    for pagina in doc:
        texto += pagina.get_text()
    return texto

def extraer_texto_pptx(ruta):
    prs = Presentation(ruta)
    texto = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto += shape.text + "\n"
    return texto

@router.post("/exam_generator")
async def exam_generator(file: UploadFile = File(...)):
    try:
        print("➡️ Inicio del endpoint /exam_generator")

        # Guardar archivo temporalmente
        file_location = f"temp_files/{file.filename}"
        os.makedirs("temp_files", exist_ok=True)
        print(f"📁 Guardando archivo en: {file_location}")

        async with aiofiles.open(file_location, 'wb') as out_file:
            content = await file.read()
            await out_file.write(content)
            await asyncio.sleep(0.1)  # pequeñoo delay para garantir que el archivo es subido antes de enviar a deepseek

        print("✅ Archivo guardado correctamente")

        # Detectar tipo de archivo y extraer texto
        if file.filename.endswith(".docx"):
            print("📄 Detectado archivo DOCX")
            texto = extraer_texto_docx(file_location)
        elif file.filename.endswith(".pdf"):
            print("📄 Detectado archivo PDF")
            texto = extraer_texto_pdf(file_location)
        elif file.filename.endswith(".pptx"):
            print("📄 Detectado archivo PPTX")
            texto = extraer_texto_pptx(file_location)
        elif file.filename.endswith(".txt"):
            print("📄 Detectado archivo TXT")
            with open(file_location, "r", encoding="utf-8") as f:
                texto = f.read()
        else:
            raise HTTPException(status_code=400, detail="Tipo de archivo no soportado. Usa .docx, .pdf, .pptx o .txt")

        print("🧠 Enviando contenido a DeepSeek")
        examen_generado = generar_examen_deepseek(texto)
        print("✅ Respuesta recibida de DeepSeek")

        os.remove(file_location)
        print("🗑️ Archivo temporal eliminado")

        return JSONResponse(content={"exam": examen_generado})

    except Exception as e:
        print(f"❌ Error detectado: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/answer_generator")
async def answer_generator(file: UploadFile = File(...), exam_text: str = Form(...)):
    try:
        print("➡️ Inicio del endpoint /answer_generator")
        file_location = f"temp_files/{file.filename}"
        os.makedirs("temp_files", exist_ok=True)

        async with aiofiles.open(file_location, 'wb') as out_file:
            content = await file.read()
            await out_file.write(content)
            await asyncio.sleep(0.1)

        # Extrai texto do arquivo (não usado no deepseek, mas pode ser usado se quiser combinar)
        if file.filename.endswith(".docx"):
            contenido = extraer_texto_docx(file_location)
        elif file.filename.endswith(".pdf"):
            contenido = extraer_texto_pdf(file_location)
        elif file.filename.endswith(".pptx"):
            contenido = extraer_texto_pptx(file_location)
        elif file.filename.endswith(".txt"):
            with open(file_location, "r", encoding="utf-8") as f:
                contenido = f.read()
        else:
            raise HTTPException(status_code=400, detail="Tipo de archivo no soportado.")

        respuestas = generar_respuestas_deepseek(contenido, exam_text)
        os.remove(file_location)
        return JSONResponse(content={"answers": respuestas})

    except Exception as e:
        print(f"❌ Error detectado: {e}")
        raise HTTPException(status_code=500, detail=str(e))